import base64
import shutil
import zipfile
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_images_pptx(path):
    prs = Presentation(path)
    images = []
    seen = set()

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                img_hash = hash(img.blob)
                if img_hash in seen:
                    continue
                seen.add(img_hash)
                ext = img.ext.lower().replace("jpeg", "jpg")
                b64 = base64.b64encode(img.blob).decode()

                # Get the internal zip path for this image
                rId = shape._element.blipFill.blip.get(
                    "{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}embed"
                ) or shape._element.blipFill.blip.attrib.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                )
                part_name = str(slide.part.rels[rId].target_part.partname).lstrip("/") if rId else None

                images.append({
                    "id": f"slide{slide_idx}_{shape.shape_id}",
                    "data": b64,
                    "ext": ext,
                    "content_type": img.content_type,
                    "_part_name": part_name,
                })

    return images


def replace_images_pptx(src, dst, selected_ids, logo_bytes, logo_ext):
    """Replace selected images using direct zip manipulation."""
    prs = Presentation(src)
    id_set = set(selected_ids)
    target_parts = set()

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                key = f"slide{slide_idx}_{shape.shape_id}"
                if key in id_set:
                    rId = shape._element.blipFill.blip.attrib.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    if rId:
                        part_name = str(slide.part.rels[rId].target_part.partname).lstrip("/")
                        target_parts.add(part_name)

    tmp = dst + ".tmp"
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename in target_parts:
                zout.writestr(item, logo_bytes)
            else:
                zout.writestr(item, data)

    os.replace(tmp, dst)
