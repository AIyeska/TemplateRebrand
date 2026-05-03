"""
Creates new branded SoftwareOne Office templates from a user prompt via Claude API.
Supports: .docx, .pptx, .xlsx
"""
import io
import json
import os
import anthropic

from docx import Document
from docx.shared import Pt, RGBColor, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PIn

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from utils.logo import get_logo_white, get_logo_dark

# SoftwareOne brand colours
SWO_BLUE   = (0,   48,  135)
SWO_CYAN   = (0,  163,  224)
SWO_DARK   = (8,   12,   20)
SWO_WHITE  = (255, 255, 255)
SWO_GREY   = (240, 242, 247)
SWO_TEXT   = (30,  30,   50)


# ─── Claude content generation ───────────────────────────────────────────────

SYSTEM_PROMPT = """Du er en ekspert på å lage profesjonelle forretningsdokumenter for MSP-bransjen.
Brukeren gir deg en beskrivelse av en mal de trenger. Du returnerer ALLTID kun gyldig JSON — ingen annen tekst.

JSON-format (tilpass innholdet til filtypen):
{
  "title": "Malnavn",
  "subtitle": "Undertittel / formål",
  "sections": [
    {
      "heading": "Seksjonstittel",
      "content": "Valgfri brødtekst",
      "items": ["Punkt 1", "Punkt 2"],
      "placeholder": "[Fyll inn her]"
    }
  ],
  "columns": ["Kolonne 1", "Kolonne 2"],
  "notes": "Valgfrie instruksjoner til brukeren"
}

Regler:
- Lag realistisk, brukbar innhold med riktige plassholdere i [hakeparenteser]
- Norsk som standard med mindre brukeren ber om engelsk
- Tilpass antall seksjoner til filtypen (Word: 4-8, PowerPoint: 5-10 slides, Excel: relevante kolonner)
- For PowerPoint: sections = slides; hver slide har heading + items (bullet-points)
- For Excel: sections = ark-seksjoner, columns = kolonneoverskrifter
- SoftwareOne-branding er allerede lagt til av appen — ikke nevn det i innholdet"""


def generate_content(prompt: str, file_type: str) -> dict:
    client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

    user_msg = f"Filtype: {file_type.upper()}\nBeskrivelse: {prompt}"

    msg = client.messages.create(
        model="claude-haiku-4-5",
        max_tokens=2048,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": user_msg}],
    )

    raw = msg.content[0].text.strip()
    # strip markdown code fences if present
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    return json.loads(raw.strip())


# ─── DOCX ────────────────────────────────────────────────────────────────────

def _rgb(r, g, b):
    return RGBColor(r, g, b)


def _set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)


def create_docx(content: dict) -> bytes:
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    # ── Logo header band ──
    logo_bytes = get_logo_white(width=280)
    logo_stream = io.BytesIO(logo_bytes)

    header_para = doc.add_paragraph()
    header_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = header_para.add_run()
    run.add_picture(logo_stream, width=Cm(7))

    # Blue rule line
    rule = doc.add_paragraph()
    rule_run = rule.add_run("─" * 80)
    rule_run.font.color.rgb = _rgb(*SWO_BLUE)
    rule_run.font.size = Pt(5)

    doc.add_paragraph()

    # ── Title ──
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    title_run = title_p.add_run(content.get("title", "Dokument"))
    title_run.font.bold  = True
    title_run.font.size  = Pt(22)
    title_run.font.color.rgb = _rgb(*SWO_BLUE)

    sub = content.get("subtitle", "")
    if sub:
        sub_p = doc.add_paragraph()
        sub_run = sub_p.add_run(sub)
        sub_run.font.size  = Pt(11)
        sub_run.font.color.rgb = _rgb(100, 110, 130)
        sub_run.font.italic = True

    doc.add_paragraph()

    # ── Sections ──
    for sec in content.get("sections", []):
        heading = doc.add_paragraph()
        h_run = heading.add_run(sec.get("heading", ""))
        h_run.font.bold  = True
        h_run.font.size  = Pt(13)
        h_run.font.color.rgb = _rgb(*SWO_BLUE)
        heading.paragraph_format.space_before = Pt(10)

        body_text = sec.get("content", "")
        if body_text:
            p = doc.add_paragraph(body_text)
            p.runs[0].font.size = Pt(10)

        placeholder = sec.get("placeholder", "")
        if placeholder:
            ph_p = doc.add_paragraph(placeholder)
            ph_run = ph_p.runs[0]
            ph_run.font.size   = Pt(10)
            ph_run.font.italic = True
            ph_run.font.color.rgb = _rgb(150, 150, 170)

        for item in sec.get("items", []):
            item_p = doc.add_paragraph(style="List Bullet")
            item_p.add_run(item).font.size = Pt(10)

    # ── Notes ──
    notes = content.get("notes", "")
    if notes:
        doc.add_paragraph()
        n_p = doc.add_paragraph()
        n_run = n_p.add_run(f"ℹ️  {notes}")
        n_run.font.size   = Pt(9)
        n_run.font.italic = True
        n_run.font.color.rgb = _rgb(100, 110, 130)

    # ── Footer ──
    footer = doc.sections[0].footer
    footer_para = footer.paragraphs[0]
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer_para.add_run("SoftwareOne Confidential  ·  www.softwareone.com")
    footer_run.font.size  = Pt(8)
    footer_run.font.color.rgb = _rgb(150, 150, 170)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ─── PPTX ────────────────────────────────────────────────────────────────────

def _prgb(r, g, b):
    return PRGBColor(r, g, b)


def _add_text_box(slide, text, l, t, w, h, font_size=18, bold=False, color=(255,255,255), align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(PIn(l), PIn(t), PIn(w), PIn(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = PPt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = _prgb(*color)
    return txBox


def _fill_shape(shape, r, g, b):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _prgb(r, g, b)


def create_pptx(content: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = PIn(13.33)
    prs.slide_height = PIn(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    slides_data = content.get("sections", [])
    title_text  = content.get("title", "Presentasjon")
    subtitle    = content.get("subtitle", "")

    # ── Title slide ──────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)

    # Dark background
    bg = slide.shapes.add_shape(1, PIn(0), PIn(0), PIn(13.33), PIn(7.5))
    _fill_shape(bg, *SWO_DARK)
    bg.line.fill.background()

    # Cyan glow line at bottom
    glow = slide.shapes.add_shape(1, PIn(0), PIn(7.1), PIn(13.33), PIn(0.05))
    _fill_shape(glow, *SWO_CYAN)
    glow.line.fill.background()

    # Orange accent line
    orange = slide.shapes.add_shape(1, PIn(0), PIn(7.15), PIn(6), PIn(0.04))
    _fill_shape(orange, 255, 106, 0)
    orange.line.fill.background()

    # Logo
    logo_bytes = get_logo_white(width=400)
    logo_stream = io.BytesIO(logo_bytes)
    prs.slides[0].shapes.add_picture(logo_stream, PIn(0.6), PIn(0.5), width=PIn(3.0))

    # Title text
    _add_text_box(slide, title_text, 0.6, 2.8, 12.0, 1.4, font_size=36, bold=True, color=SWO_WHITE)
    if subtitle:
        _add_text_box(slide, subtitle, 0.6, 4.0, 10.0, 0.8, font_size=16, color=(180, 190, 210))

    # ── Content slides ────────────────────────────────────
    for sec in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # White background
        bg = slide.shapes.add_shape(1, PIn(0), PIn(0), PIn(13.33), PIn(7.5))
        _fill_shape(bg, 255, 255, 255)
        bg.line.fill.background()

        # Blue header bar
        header = slide.shapes.add_shape(1, PIn(0), PIn(0), PIn(13.33), PIn(1.15))
        _fill_shape(header, *SWO_BLUE)
        header.line.fill.background()

        # Cyan accent stripe
        stripe = slide.shapes.add_shape(1, PIn(0), PIn(1.15), PIn(13.33), PIn(0.05))
        _fill_shape(stripe, *SWO_CYAN)
        stripe.line.fill.background()

        # Slide heading in header
        heading = sec.get("heading", "")
        _add_text_box(slide, heading, 0.35, 0.2, 10.5, 0.8, font_size=20, bold=True, color=SWO_WHITE)

        # Logo small in header (right)
        logo_stream = io.BytesIO(get_logo_white(width=200))
        slide.shapes.add_picture(logo_stream, PIn(11.2), PIn(0.18), width=PIn(1.85))

        # Body content
        y = 1.5
        body = sec.get("content", "")
        if body:
            _add_text_box(slide, body, 0.5, y, 12.0, 0.6, font_size=12, color=SWO_TEXT)
            y += 0.7

        placeholder = sec.get("placeholder", "")
        if placeholder:
            _add_text_box(slide, placeholder, 0.5, y, 12.0, 0.5, font_size=11, color=(150, 155, 175))
            y += 0.6

        items = sec.get("items", [])
        for item in items:
            _add_text_box(slide, f"  •  {item}", 0.5, y, 12.0, 0.45, font_size=12, color=SWO_TEXT)
            y += 0.45
            if y > 6.8:
                break

        # Footer
        footer_bar = slide.shapes.add_shape(1, PIn(0), PIn(7.25), PIn(13.33), PIn(0.25))
        _fill_shape(footer_bar, *SWO_BLUE)
        footer_bar.line.fill.background()
        _add_text_box(slide, "SoftwareOne Confidential  ·  www.softwareone.com",
                      0.3, 7.27, 12.0, 0.22, font_size=7.5, color=(180, 200, 230))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── XLSX ────────────────────────────────────────────────────────────────────

def _hex(r, g, b):
    return f"{r:02X}{g:02X}{b:02X}"


def create_xlsx(content: dict) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = content.get("title", "Mal")[:31]

    blue_fill  = PatternFill("solid", fgColor=_hex(*SWO_BLUE))
    cyan_fill  = PatternFill("solid", fgColor=_hex(*SWO_CYAN))
    grey_fill  = PatternFill("solid", fgColor="F0F2F7")
    white_fill = PatternFill("solid", fgColor="FFFFFF")

    thin = Side(border_style="thin", color="C0CCDF")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    # ── Logo row (row 1) ──
    ws.row_dimensions[1].height = 45
    ws.merge_cells("A1:H1")
    logo_cell = ws["A1"]
    logo_cell.value = "  SoftwareOne"
    logo_cell.font  = Font(name="Segoe UI", bold=True, size=18, color="FFFFFF")
    logo_cell.fill  = blue_fill
    logo_cell.alignment = Alignment(horizontal="left", vertical="center")

    # ── Title row (row 2) ──
    ws.row_dimensions[2].height = 30
    ws.merge_cells("A2:H2")
    title_cell = ws["A2"]
    title_cell.value = content.get("title", "")
    title_cell.font  = Font(name="Segoe UI", bold=True, size=14, color="FFFFFF")
    title_cell.fill  = PatternFill("solid", fgColor="002080")
    title_cell.alignment = Alignment(horizontal="left", vertical="center", indent=1)

    # ── Subtitle row (row 3) ──
    sub = content.get("subtitle", "")
    ws.row_dimensions[3].height = 20
    ws.merge_cells("A3:H3")
    sub_cell = ws["A3"]
    sub_cell.value = sub
    sub_cell.font  = Font(name="Segoe UI", italic=True, size=10, color="556080")
    sub_cell.fill  = grey_fill
    sub_cell.alignment = Alignment(horizontal="left", vertical="center", indent=1)

    ws.row_dimensions[4].height = 8  # spacer

    # ── Column headers (row 5) ──
    columns = content.get("columns", [])
    if not columns:
        columns = ["Felt", "Verdi", "Kommentar"]

    ws.row_dimensions[5].height = 22
    for ci, col in enumerate(columns, start=1):
        cell = ws.cell(row=5, column=ci, value=col)
        cell.font      = Font(name="Segoe UI", bold=True, size=10, color="FFFFFF")
        cell.fill      = cyan_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border    = border
        ws.column_dimensions[get_column_letter(ci)].width = max(18, len(col) + 4)

    # ── Data rows from sections ──
    row = 6
    for sec in content.get("sections", []):
        # Section heading row
        ws.row_dimensions[row].height = 20
        ws.merge_cells(f"A{row}:{get_column_letter(max(len(columns),1))}{row}")
        h_cell = ws.cell(row=row, column=1, value=sec.get("heading", ""))
        h_cell.font      = Font(name="Segoe UI", bold=True, size=10, color=_hex(*SWO_BLUE))
        h_cell.fill      = grey_fill
        h_cell.alignment = Alignment(horizontal="left", vertical="center", indent=1)
        row += 1

        items = sec.get("items", [])
        placeholder = sec.get("placeholder", "")

        if items:
            for item in items:
                ws.row_dimensions[row].height = 18
                cell = ws.cell(row=row, column=1, value=item)
                cell.font      = Font(name="Segoe UI", size=10)
                cell.fill      = white_fill
                cell.alignment = Alignment(horizontal="left", vertical="center", indent=2)
                cell.border    = border
                # Empty cells for remaining columns
                for ci in range(2, len(columns) + 1):
                    ph_cell = ws.cell(row=row, column=ci, value=placeholder if ci == 2 else "")
                    ph_cell.font      = Font(name="Segoe UI", size=10, italic=True, color="AAAAAA")
                    ph_cell.fill      = white_fill
                    ph_cell.border    = border
                row += 1
        else:
            ws.row_dimensions[row].height = 18
            ws.cell(row=row, column=1, value=placeholder or "[Legg til data]").font = Font(
                name="Segoe UI", size=10, italic=True, color="AAAAAA")
            row += 1

        row += 1  # blank row between sections

    # ── Notes ──
    notes = content.get("notes", "")
    if notes:
        row += 1
        ws.merge_cells(f"A{row}:H{row}")
        n_cell = ws.cell(row=row, column=1, value=f"ℹ  {notes}")
        n_cell.font      = Font(name="Segoe UI", italic=True, size=9, color="666688")
        n_cell.alignment = Alignment(horizontal="left", vertical="center", indent=1)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ─── Dispatcher ──────────────────────────────────────────────────────────────

EXT_MAP = {
    "docx": ("docx", create_docx),
    "pptx": ("pptx", create_pptx),
    "xlsx": ("xlsx", create_xlsx),
}


def create_template(prompt: str, file_type: str) -> tuple[bytes, str]:
    """Returns (file_bytes, extension)."""
    ft = file_type.lower().strip(".")
    if ft not in EXT_MAP:
        raise ValueError(f"Ikke støttet filtype: {ft}")
    ext, creator_fn = EXT_MAP[ft]
    content = generate_content(prompt, ext)
    file_bytes = creator_fn(content)
    return file_bytes, ext
